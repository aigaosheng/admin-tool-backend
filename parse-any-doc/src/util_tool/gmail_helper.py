import os.path
import base64
import email
from datetime import datetime, timezone
from google.auth.transport.requests import Request
from google.oauth2.credentials import Credentials
from google_auth_oauthlib.flow import InstalledAppFlow
from googleapiclient.discovery import build
import mimetypes
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
from email.mime.base import MIMEBase
from email import encoders
import os, io, zipfile
from googleapiclient.errors import HttpError
import PyPDF2
from docx import Document as DocxDocument
from pptx import Presentation
from PIL import Image
from pydub import AudioSegment

import pandas as pd

SCOPES = ["https://www.googleapis.com/auth/gmail.readonly",
          "https://www.googleapis.com/auth/gmail.send"
        ]

OUTPUT_DIR = "/home/gs/work/recruitAI/output"
SAVE_LOCAL = False

def save_attachment(part, msg_id, store_dir=OUTPUT_DIR):
    filename = part.get_filename()
    if not filename:
        return None

    os.makedirs(store_dir, exist_ok=True)
    path = os.path.join(store_dir, f"{msg_id}_{filename}")

    data = part.get_payload(decode=True)
    if data:
        with open(path, 'wb') as f:
            f.write(data)
        print(f"Saved attachment: {path}")
        return path
    return None


def process_attachment(file_path):
    ext = os.path.splitext(file_path)[1].lower()

    try:
        if ext == ".pdf":
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                text = ''
                for page in reader.pages:
                    text += page.extract_text()
                return {"type": "pdf", "text": text[:500] + "..."}
        elif ext == ".docx":
            doc = DocxDocument(file_path)
            text = "\n".join([para.text for para in doc.paragraphs])
            return {"type": "docx", "text": text[:500] + "..."}
        elif ext == ".pptx":
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return {"type": "pptx", "text": text[:500] + "..."}

        elif ext in [".jpg", ".jpeg", ".png"]:
            img = Image.open(file_path)
            return {"type": "image", "size": img.size, "mode": img.mode}
        elif ext == ".mp3":
            audio = AudioSegment.from_mp3(file_path)
            duration = round(len(audio) / 1000, 2)
            return {"type": "audio", "duration_sec": duration}
        else:
            return {"type": "unknown", "file": file_path}
    except Exception as e:
        return {"type": "error", "error": str(e)}
    
def extract_html_from_raw(msg):
    # raw_bytes = base64.urlsafe_b64decode(raw_encoded.encode('ASCII'))
    # msg = email.message_from_bytes(raw_bytes)
    res = "<p></p>"
    if msg.is_multipart():
        for part in msg.walk():
            content_type = part.get_content_type()
            if content_type == "text/html":
                charset = part.get_content_charset() or 'utf-8'
                res = part.get_payload(decode=True).decode(charset)
    else:
        if msg.get_content_type() == "text/html":
            charset = msg.get_content_charset() or 'utf-8'
            res = msg.get_payload(decode=True).decode(charset)
    body = {"text": "", "html": res}
    return body

def get_gmail_service(user_id):
    token_file = f"tokens/{user_id}_token.json"
    # print(f"*** token_file -> {token_file}")
    creds = None
    if os.path.exists(token_file):
        creds = Credentials.from_authorized_user_file(token_file, SCOPES)
    if not creds or not creds.valid:
        if creds and creds.expired and creds.refresh_token:
            creds.refresh(Request())
        else:
            flow = InstalledAppFlow.from_client_secrets_file("credentials.json", SCOPES)
            creds = flow.run_local_server(port=0)
        os.makedirs("tokens", exist_ok=True)
        with open(token_file, "w") as token:
            token.write(creds.to_json())
    return build("gmail", "v1", credentials=creds)

def list_messages_v1(service, max_results=50, query="is:unread"):
    """Return list of Gmail messages (dicts compatible with existing inbox).
    labelIds=['INBOX'], q='is:unread')
    """
    results = (
        service.users()
        .messages()
        .list(userId="me", maxResults=max_results, labelIds=["INBOX"], q=query)
        .execute()
    )
    messages = results.get("messages", [])
    parsed = []
    for msg in messages:
        txt = (
            service.users()
            .messages()
            .get(userId="me", id=msg["id"], format="full")
            .execute()
        )
        headers = {h["name"]: h["value"] for h in txt["payload"]["headers"]}
        snippet = txt.get("snippet", "")
        subject = headers.get("Subject", "(no subject)")
        sender = headers.get("From", "")
        date_str = headers.get("Date", "")
        # try:
        #     date = email.utils.parsedate_to_datetime(date_str).replace(tzinfo=timezone.utc)
        # except Exception:
        #     date = datetime.now(timezone.utc)
        try:
            date = pd.to_datetime(date_str, utc=True).tz_localize(None)  # naive
        except Exception:
            date = pd.Timestamp.utcnow().tz_localize(None)

        parsed.append(
            {
                "id": msg["id"],
                "from": sender,
                "subject": subject,
                "body": snippet,
                "date": date,
                "read": "UNREAD" not in txt.get("labelIds", []),
                "priority": "High" if "IMPORTANT" in txt.get("labelIds", []) else "Medium",  
                "category": "External",  # Simplistic mapping
            }
        )
    
        # print(f"*** email ** => {parsed}")
    return parsed

def extract_body(msg):
    text_content = ""
    html_content = ""

    def parse_part(part):
        nonlocal text_content, html_content
        content_type = part.get_content_type()
        payload = part.get_payload(decode=True)

        if payload is None:
            return

        charset = part.get_content_charset() or 'utf-8'
        decoded = payload.decode(charset, errors="ignore")

        if content_type == "text/plain":
            text_content += decoded
        elif content_type == "text/html":
            html_content += decoded

    if msg.is_multipart():
        for part in msg.walk():
            parse_part(part)
    else:
        parse_part(msg)

    return {'text': text_content.strip(), 'html': html_content.strip()}

def list_messages(service, max_results=50, query="is:unread"):
    """Return list of Gmail messages (dicts compatible with existing inbox).
    labelIds=['INBOX'], q='is:unread')
    """
    """Return list of Gmail messages (dicts compatible with existing inbox).
    labelIds=['INBOX'], q='is:unread')
    """
    results = (
        service.users()
        .messages()
        .list(userId="me", maxResults=max_results, labelIds=["INBOX"], q=query)
        .execute()
    )
    messages = results.get("messages", [])
    parsed = []
    for msg in messages:
        message = (
            service.users()
            .messages()
            .get(userId="me", id=msg["id"], format="raw")
            .execute()
        )
        msg_raw = base64.urlsafe_b64decode(message['raw'].encode('ASCII'))
        msg_data = email.message_from_bytes(msg_raw, policy=email.policy.default)

        # print(f"*** msg_data -> {msg_data.keys()}")

        subject = msg_data.get("Subject", "(no subject)")
        sender = msg_data.get("From", "")
        date_str = msg_data.get("Date", "")
        receiver = msg_data.get('To', "")

        try:
            date = pd.to_datetime(date_str, utc=True).tz_localize(None)  # naive
        except Exception:
            date = pd.Timestamp.utcnow().tz_localize(None)

        # body = extract_body(msg_data)
        body = extract_html_from_raw(msg_data)
    
        # print(f"*** body -> {body}")

        attachments_info = []

        # Process parts
        if msg_data.is_multipart():
            for part in msg_data.walk():
                if part.get_filename():
                    if SAVE_LOCAL:
                        file_path = save_attachment(part, msg["id"])
                        if file_path:
                            info = process_attachment(file_path)
                            info['filename'] = part.get_filename()
                            attachments_info.append(info)
                    else:
                        # print(f"*** part ==> {dir(part)}")
                        # print(f"*** part['filename'] => {part.get_filename()}")
                        # print(f"*** part['mime'] => {part.get_content_type()}")
                        data = part.get_payload(decode=True)                        
                        attachments_info.append({
                            'filename': part.get_filename(),
                            'mime': part.get_content_type(),
                            'data': data
                        })

        parsed.append(
            {
                "id": msg["id"],
                "from": sender,
                "subject": subject,
                "body": body.get("text", "").strip('\n'),
                "body_html": body.get("html", "").strip(),
                "date": date,
                "read": "UNREAD" not in message.get("labelIds", []),
                "priority": "High" if "IMPORTANT" in message.get("labelIds", []) else "Medium",  
                "category": "External",  # Simplistic mapping
                "attachments": attachments_info
            }
        )
    

    return parsed

def send_gmail(service, to, subject, body):
    message = MIMEText(body)
    message["to"] = to
    message["subject"] = subject
    raw = base64.urlsafe_b64encode(message.as_bytes()).decode()
    try:
        sd_email = service.users().messages().send(userId="me", body={"raw": raw}).execute()
        return sd_email
    except Exception as e:
        print(f"** e ** {e}** ")
        return None

def create_message_with_attachments(to, subject, body, file_list):
    """
    file_list: list of (filename, bytes) tuples
    returns base64url-safe string ready for the Gmail API
    """
    msg = MIMEMultipart()
    msg["to"] = to
    msg["subject"] = subject

    msg.attach(MIMEText(body, "plain"))

    for name, data in file_list:
        main, sub = mimetypes.guess_type(name)[0].split("/", 1)
        part = MIMEBase(main, sub)
        part.set_payload(data)
        encoders.encode_base64(part)
        part.add_header(
            "Content-Disposition",
            f'attachment; filename="{name}"',
        )
        msg.attach(part)

    raw = base64.urlsafe_b64encode(msg.as_bytes()).decode()
    return {"raw": raw}

def send_gmail_with_attachments(service, to, subject, body, file_list):
    message = create_message_with_attachments(to, subject, body, file_list)
    return (
        service.users()
        .messages()
        .send(userId="me", body=message)
        .execute()
    )

def get_attachments(service, msg_id):
    """
    Returns a list of dicts:
    [
      {'filename': 'resume.pdf', 'mime': 'application/pdf', 'data': b'...'},
      ...
    ]
    """
    try:
        message = service.users().messages().get(
            userId='me', id=msg_id, format='full'
        ).execute()
    except HttpError:
        return []

    attachments = []
    for part in message.get('payload', {}).get('parts', []):
        if part.get('filename') and part.get('body', {}).get('attachmentId'):
            att_id = part['body']['attachmentId']
            att = service.users().messages().attachments().get(
                userId='me', messageId=msg_id, id=att_id
            ).execute()
            data = base64.urlsafe_b64decode(att['data'])
            attachments.append({
                'filename': part['filename'],
                'mime': part['mimeType'],
                'data': data
            })
    return attachments